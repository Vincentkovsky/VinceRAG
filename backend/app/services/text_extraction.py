"""
Text extraction service for various document formats
"""

import io
import csv
import hashlib
from typing import Dict, Any, Optional, Tuple
from pathlib import Path

# PDF processing
import pypdf
import pdfplumber

# Office documents
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# File type detection
import magic

# Import removed - not needed for this service


class TextExtractionError(Exception):
    """Custom exception for text extraction errors"""
    pass


class TextExtractor:
    """Service for extracting text from various document formats"""
    
    # Maximum file size (50MB)
    MAX_FILE_SIZE = 50 * 1024 * 1024
    
    # Supported file types
    SUPPORTED_TYPES = {
        'pdf', 'docx', 'txt', 'md', 'pptx', 'xlsx', 'csv', 'rtf'
    }
    
    def __init__(self):
        self.mime_detector = magic.Magic(mime=True)
    
    def extract_text(self, file_path: str, file_type: str = None) -> Dict[str, Any]:
        """
        Extract text from a file
        
        Args:
            file_path: Path to the file
            file_type: Optional file type hint
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise TextExtractionError(f"File not found: {file_path}")
        
        # Check file size
        file_size = file_path.stat().st_size
        if file_size > self.MAX_FILE_SIZE:
            raise TextExtractionError(f"File too large: {file_size} bytes (max: {self.MAX_FILE_SIZE})")
        
        # Detect file type if not provided
        if not file_type:
            file_type = self._detect_file_type(file_path)
        
        if file_type not in self.SUPPORTED_TYPES:
            raise TextExtractionError(f"Unsupported file type: {file_type}")
        
        # Calculate file hash for duplicate detection
        file_hash = self._calculate_file_hash(file_path)
        
        # Extract text based on file type
        extraction_methods = {
            'pdf': self._extract_pdf,
            'docx': self._extract_docx,
            'txt': self._extract_text_file,
            'md': self._extract_text_file,
            'rtf': self._extract_text_file,
            'pptx': self._extract_pptx,
            'xlsx': self._extract_xlsx,
            'csv': self._extract_csv
        }
        
        try:
            text_content, metadata = extraction_methods[file_type](file_path)
            
            return {
                'text': text_content,
                'file_hash': file_hash,
                'file_size': file_size,
                'file_type': file_type,
                'metadata': metadata
            }
        except Exception as e:
            raise TextExtractionError(f"Failed to extract text from {file_type} file: {str(e)}")
    
    def _detect_file_type(self, file_path: Path) -> str:
        """Detect file type using magic numbers and extension"""
        try:
            mime_type = self.mime_detector.from_file(str(file_path))
            
            # Map MIME types to our file types
            mime_map = {
                'application/pdf': 'pdf',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
                'text/plain': 'txt',
                'text/csv': 'csv',
                'application/rtf': 'rtf'
            }
            
            if mime_type in mime_map:
                return mime_map[mime_type]
            
            # Fallback to extension
            extension = file_path.suffix.lower().lstrip('.')
            if extension in self.SUPPORTED_TYPES:
                return extension
            
            # Default to text for unknown types
            return 'txt'
            
        except Exception:
            # Fallback to extension if magic fails
            extension = file_path.suffix.lower().lstrip('.')
            return extension if extension in self.SUPPORTED_TYPES else 'txt'
    
    def _calculate_file_hash(self, file_path: Path) -> str:
        """Calculate SHA-256 hash of file for duplicate detection"""
        hash_sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_sha256.update(chunk)
        return hash_sha256.hexdigest()
    
    def _extract_pdf(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PDF files using PyPDF2 and pdfplumber"""
        text_content = []
        metadata = {}
        
        try:
            # First try with pdfplumber for better text extraction
            with pdfplumber.open(file_path) as pdf:
                metadata.update({
                    'page_count': len(pdf.pages),
                    'extraction_method': 'pdfplumber'
                })
                
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(page_text)
                
                # Extract document metadata if available
                if pdf.metadata:
                    metadata.update({
                        'title': pdf.metadata.get('Title'),
                        'author': pdf.metadata.get('Author'),
                        'subject': pdf.metadata.get('Subject'),
                        'creator': pdf.metadata.get('Creator'),
                        'producer': pdf.metadata.get('Producer'),
                        'creation_date': pdf.metadata.get('CreationDate'),
                        'modification_date': pdf.metadata.get('ModDate')
                    })
        
        except Exception as e:
            # Fallback to pypdf
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = pypdf.PdfReader(file)
                    metadata.update({
                        'page_count': len(pdf_reader.pages),
                        'extraction_method': 'pypdf',
                        'fallback_reason': str(e)
                    })
                    
                    for page in pdf_reader.pages:
                        text_content.append(page.extract_text())
                    
                    # Extract metadata
                    if pdf_reader.metadata:
                        metadata.update({
                            'title': pdf_reader.metadata.get('/Title'),
                            'author': pdf_reader.metadata.get('/Author'),
                            'subject': pdf_reader.metadata.get('/Subject'),
                            'creator': pdf_reader.metadata.get('/Creator'),
                            'producer': pdf_reader.metadata.get('/Producer')
                        })
            except Exception as fallback_error:
                raise TextExtractionError(f"Both pdfplumber and pypdf failed: {fallback_error}")
        
        return '\n\n'.join(text_content), metadata
    
    def _extract_docx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from DOCX files"""
        doc = DocxDocument(file_path)
        
        # Extract text from paragraphs
        paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        
        # Extract text from tables
        table_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    table_text.append(' | '.join(row_text))
        
        # Combine all text
        all_text = paragraphs + table_text
        
        # Extract metadata
        metadata = {
            'paragraph_count': len(paragraphs),
            'table_count': len(doc.tables),
            'extraction_method': 'python-docx'
        }
        
        # Document properties
        if hasattr(doc.core_properties, 'title') and doc.core_properties.title:
            metadata['title'] = doc.core_properties.title
        if hasattr(doc.core_properties, 'author') and doc.core_properties.author:
            metadata['author'] = doc.core_properties.author
        if hasattr(doc.core_properties, 'subject') and doc.core_properties.subject:
            metadata['subject'] = doc.core_properties.subject
        if hasattr(doc.core_properties, 'created') and doc.core_properties.created:
            metadata['created'] = doc.core_properties.created.isoformat()
        if hasattr(doc.core_properties, 'modified') and doc.core_properties.modified:
            metadata['modified'] = doc.core_properties.modified.isoformat()
        
        return '\n\n'.join(all_text), metadata
    
    def _extract_pptx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PowerPoint files"""
        prs = Presentation(file_path)
        
        slide_texts = []
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                slide_texts.append(f"Slide {slide_num + 1}:\n" + '\n'.join(slide_text))
        
        metadata = {
            'slide_count': len(prs.slides),
            'extraction_method': 'python-pptx'
        }
        
        # Extract presentation properties
        if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
            metadata['title'] = prs.core_properties.title
        if hasattr(prs.core_properties, 'author') and prs.core_properties.author:
            metadata['author'] = prs.core_properties.author
        if hasattr(prs.core_properties, 'subject') and prs.core_properties.subject:
            metadata['subject'] = prs.core_properties.subject
        
        return '\n\n'.join(slide_texts), metadata
    
    def _extract_xlsx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from Excel files"""
        workbook = load_workbook(file_path, data_only=True)
        
        sheet_texts = []
        total_rows = 0
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = []
            
            # Extract text from cells
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell_value in row:
                    if cell_value is not None:
                        row_text.append(str(cell_value).strip())
                
                if row_text:
                    sheet_text.append(' | '.join(row_text))
                    total_rows += 1
            
            if sheet_text:
                sheet_texts.append(f"Sheet: {sheet_name}\n" + '\n'.join(sheet_text))
        
        metadata = {
            'sheet_count': len(workbook.sheetnames),
            'total_rows': total_rows,
            'sheet_names': workbook.sheetnames,
            'extraction_method': 'openpyxl'
        }
        
        # Extract workbook properties
        if hasattr(workbook.properties, 'title') and workbook.properties.title:
            metadata['title'] = workbook.properties.title
        if hasattr(workbook.properties, 'creator') and workbook.properties.creator:
            metadata['author'] = workbook.properties.creator
        if hasattr(workbook.properties, 'subject') and workbook.properties.subject:
            metadata['subject'] = workbook.properties.subject
        
        return '\n\n'.join(sheet_texts), metadata
    
    def _extract_csv(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from CSV files"""
        rows = []
        row_count = 0
        
        # Try different encodings
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding, newline='') as csvfile:
                    # Detect delimiter
                    sample = csvfile.read(1024)
                    csvfile.seek(0)
                    sniffer = csv.Sniffer()
                    delimiter = sniffer.sniff(sample).delimiter
                    
                    reader = csv.reader(csvfile, delimiter=delimiter)
                    
                    for row in reader:
                        if any(cell.strip() for cell in row):  # Skip empty rows
                            rows.append(' | '.join(cell.strip() for cell in row))
                            row_count += 1
                
                break  # Success, exit encoding loop
                
            except (UnicodeDecodeError, UnicodeError):
                continue  # Try next encoding
            except Exception as e:
                if encoding == encodings[-1]:  # Last encoding failed
                    raise TextExtractionError(f"Failed to read CSV file: {e}")
        
        metadata = {
            'row_count': row_count,
            'encoding': encoding,
            'delimiter': delimiter,
            'extraction_method': 'csv'
        }
        
        return '\n'.join(rows), metadata
    
    def _extract_text_file(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from plain text files (TXT, MD, RTF)"""
        # Try different encodings
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    content = file.read()
                
                metadata = {
                    'encoding': encoding,
                    'line_count': content.count('\n') + 1,
                    'character_count': len(content),
                    'extraction_method': 'text'
                }
                
                return content, metadata
                
            except (UnicodeDecodeError, UnicodeError):
                continue  # Try next encoding
            except Exception as e:
                if encoding == encodings[-1]:  # Last encoding failed
                    raise TextExtractionError(f"Failed to read text file: {e}")
        
        raise TextExtractionError("Could not decode text file with any supported encoding")